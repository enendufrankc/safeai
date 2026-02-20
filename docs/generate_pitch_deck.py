"""Generate SafeAI investor pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x16, 0xC7, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
SUBTLE_TEXT = RGBColor(0x6C, 0x75, 0x7D)
CORAL = RGBColor(0xE7, 0x4C, 0x3C)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
DARK_BG = RGBColor(0x0F, 0x0F, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=24,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=20,
                   color=DARK_TEXT, spacing=Pt(8), bold=False, font_name="Calibri",
                   alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = spacing
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.6), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Decorative teal accent line at top
add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.05), TEAL)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "SafeAI", 72, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.0),
             "The Open-Source Privacy & Security Control Layer\nfor AI Agents",
             28, TEAL, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Framework-agnostic  \u2022  Deterministic  \u2022  Lightweight  \u2022  Open Source",
             18, RGBColor(0x99, 0x99, 0xAA), False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "The Problem", 44, NAVY, True)

add_text_box(slide, Inches(1.0), Inches(1.9), Inches(10), Inches(0.6),
             "AI agents are powerful but unsafe by default.", 24, CORAL, True)

# Problem cards
problems = [
    ("Agents Leak Data", "Public incidents: platforms exposing\nuser data, credentials, API keys"),
    ("Current Fixes Are Fragile", "Prompt-based safety drifts.\nFramework guardrails don't transfer."),
    ("The Gap Is Growing", "Agent deployments accelerating.\nRegulation tightening. Controls falling behind."),
]

for i, (title, desc) in enumerate(problems):
    x = Inches(1.0 + i * 3.8)
    y = Inches(3.0)
    card = add_shape_bg(slide, x, y, Inches(3.4), Inches(3.2), LIGHT_GRAY)
    card.shadow.inherit = False
    add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(2.8), Inches(0.6),
                 title, 22, NAVY, True)
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.2), Inches(2.8), Inches(1.8),
                   desc.split("\n"), 16, SUBTLE_TEXT)

add_text_box(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(0.5),
             "Major tech companies have banned agent platforms over security concerns. The problem demands an industry-level solution.",
             16, SUBTLE_TEXT, False, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "The Solution", 44, NAVY, True)

add_text_box(slide, Inches(1.0), Inches(1.9), Inches(11), Inches(0.8),
             "A runtime control layer that enforces privacy and security policies\nat the boundaries where AI agent data flows.",
             22, DARK_TEXT, False)

# Three boundary boxes
boundaries = [
    ("Input Boundary", "Classify & filter data\nbefore the AI processes it", "\u2B9E"),
    ("Action Boundary", "Validate tool calls, inject\nscoped credentials, gate actions", "\u2699"),
    ("Output Boundary", "Scan responses, redact\nsensitive data, block violations", "\u2B9C"),
]

for i, (title, desc, icon) in enumerate(boundaries):
    x = Inches(1.0 + i * 3.8)
    y = Inches(3.2)
    card = add_shape_bg(slide, x, y, Inches(3.4), Inches(2.8), NAVY)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.8), Inches(0.5),
                 icon, 32, TEAL, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(2.8), Inches(0.5),
                 title, 20, WHITE, True, PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.5), Inches(2.8), Inches(1.2),
                   desc.split("\n"), 15, RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.0), Inches(6.4), Inches(11), Inches(0.6),
             "The agent cannot leak what it never received.",
             20, TEAL, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: How It Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(1.0), Inches(0.6), Inches(10), Inches(0.8),
             "How It Works", 44, WHITE, True)

add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.6),
             "SafeAI sits underneath any AI framework \u2014 not inside it", 20, TEAL)

# Flow diagram
flow_items = [
    ("User Input", SUBTLE_TEXT),
    ("\u2192", TEAL),
    ("Input\nScanner", TEAL),
    ("\u2192", TEAL),
    ("Agent\nReasoning", SUBTLE_TEXT),
    ("\u2192", TEAL),
    ("Action\nInterceptor", TEAL),
    ("\u2192", TEAL),
    ("Output\nGuard", TEAL),
    ("\u2192", TEAL),
    ("Safe\nResponse", SUBTLE_TEXT),
]

x_start = Inches(0.5)
for i, (text, color) in enumerate(flow_items):
    x = x_start + Inches(i * 1.15)
    if "\u2192" in text:
        add_text_box(slide, x, Inches(3.2), Inches(0.8), Inches(0.8),
                     text, 28, color, True, PP_ALIGN.CENTER)
    elif color == TEAL:
        card = add_shape_bg(slide, x, Inches(2.8), Inches(1.1), Inches(1.4),
                           RGBColor(0x25, 0x25, 0x3E))
        add_text_box(slide, x + Inches(0.05), Inches(3.0), Inches(1.0), Inches(1.0),
                     text, 14, TEAL, True, PP_ALIGN.CENTER)
    else:
        add_text_box(slide, x, Inches(3.0), Inches(1.1), Inches(1.0),
                     text, 14, RGBColor(0x88, 0x88, 0x99), False, PP_ALIGN.CENTER)

# Bottom features
features = [
    "Policy Engine", "Data Classifier", "Memory Controller", "Audit Logger"
]
for i, feat in enumerate(features):
    x = Inches(1.5 + i * 2.8)
    add_shape_bg(slide, x, Inches(5.2), Inches(2.2), Inches(0.7),
                RGBColor(0x25, 0x25, 0x3E))
    add_text_box(slide, x, Inches(5.3), Inches(2.2), Inches(0.5),
                 feat, 16, WHITE, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
             "< 50ms overhead  \u2022  No extra LLM calls  \u2022  Deterministic rules  \u2022  Full audit trail",
             16, RGBColor(0x88, 0x88, 0x99), False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Key Differentiators
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "Why SafeAI", 44, NAVY, True)

diffs = [
    ("Framework-Agnostic", "Works with LangChain, Claude ADK,\nGoogle ADK, custom stacks"),
    ("Open Source", "Transparent, inspectable,\ncommunity-driven"),
    ("Boundary-Based", "Enforces at data boundaries,\nnot inside prompts"),
    ("Deterministic", "Rules-based, predictable\nbehavior every time"),
    ("Lightweight", "< 50ms overhead,\nno extra LLM calls"),
    ("Developer-First", "SDK, CLI, YAML config \u2014\nnot a heavy platform"),
]

for i, (title, desc) in enumerate(diffs):
    col = i % 3
    row = i // 3
    x = Inches(1.0 + col * 3.8)
    y = Inches(2.2 + row * 2.5)

    bar = add_shape_bg(slide, x, y, Inches(0.06), Inches(1.8), TEAL)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(3.0), Inches(0.5),
                 title, 20, NAVY, True)
    add_multi_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.0), Inches(1.0),
                   desc.split("\n"), 15, SUBTLE_TEXT)

# ============================================================
# SLIDE 6: Market Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "Market Opportunity", 44, NAVY, True)

# TAM/SAM/SOM circles (represented as rounded rectangles)
market_data = [
    ("TAM", "$10B+", "AI Infrastructure\nSecurity by 2028"),
    ("SAM", "$2\u20133B", "AI Agent Security\nSpecifically"),
    ("SOM", "$5M", "Year 3 Target\n100+ Enterprise Customers"),
]

for i, (label, amount, desc) in enumerate(market_data):
    x = Inches(1.0 + i * 3.8)
    y = Inches(2.2)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), y, Inches(2.4), Inches(2.4))
    circle.fill.solid()
    colors = [NAVY, RGBColor(0x25, 0x25, 0x3E), RGBColor(0x10, 0x80, 0x6B)]
    circle.fill.fore_color.rgb = colors[i]
    circle.line.fill.background()

    add_text_box(slide, x + Inches(0.5), y + Inches(0.4), Inches(2.4), Inches(0.4),
                 label, 18, TEAL, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.8), Inches(2.4), Inches(0.6),
                 amount, 36, WHITE, True, PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.5), y + Inches(1.5), Inches(2.4), Inches(0.8),
                   desc.split("\n"), 13, RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

# Why now
add_text_box(slide, Inches(1.0), Inches(5.0), Inches(10), Inches(0.5),
             "Why Now", 24, NAVY, True)

why_now = [
    "Public incidents creating urgency \u2014 companies can no longer defer AI security",
    "No open standard exists \u2014 first mover establishes the category",
    "Agent adoption accelerating \u2014 every deployment needs security controls",
    "Regulation arriving \u2014 EU AI Act, GDPR enforcement creating compliance requirements",
]
add_multi_text(slide, Inches(1.0), Inches(5.5), Inches(10), Inches(1.8),
               why_now, 15, SUBTLE_TEXT, Pt(6))

# ============================================================
# SLIDE 7: Business Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "Business Model: Open-Core", 44, NAVY, True)

tiers = [
    ("Open Source", "Free Forever",
     ["Core engine & SDK", "Policy engine", "CLI tools", "Framework adapters",
      "Audit logging", "Community support"],
     LIGHT_GRAY, NAVY),
    ("SafeAI Pro", "$500\u2013$2K/mo",
     ["Web dashboard", "Approval workflow UI", "Compliance reports",
      "Multi-tenant policies", "Priority support", "SLA guarantees"],
     NAVY, WHITE),
    ("SafeAI Enterprise", "Custom",
     ["Managed gateway", "SSO & RBAC", "Dedicated support",
      "Custom classifiers", "On-premise deploy", "Security audit docs"],
     RGBColor(0x10, 0x80, 0x6B), WHITE),
]

for i, (name, price, features, bg_color, text_color) in enumerate(tiers):
    x = Inches(1.0 + i * 3.8)
    y = Inches(2.0)
    card = add_shape_bg(slide, x, y, Inches(3.4), Inches(4.8), bg_color)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.8), Inches(0.5),
                 name, 22, text_color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(2.8), Inches(0.5),
                 price, 20, TEAL if bg_color != LIGHT_GRAY else NAVY, True, PP_ALIGN.CENTER)

    feat_color = SUBTLE_TEXT if bg_color == LIGHT_GRAY else RGBColor(0xBB, 0xBB, 0xCC)
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.6), Inches(2.8), Inches(3.0),
                   [f"\u2713  {f}" for f in features], 14, feat_color, Pt(6))

# ============================================================
# SLIDE 8: Competitive Landscape
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "Competitive Advantage", 44, NAVY, True)

# Table-like comparison
headers = ["", "Framework\nGuardrails", "Enterprise\nPlatforms", "DLP\nTools", "SafeAI"]
header_y = Inches(2.2)
col_w = Inches(2.2)
row_h = Inches(0.65)

# Header row
for i, h in enumerate(headers):
    x = Inches(1.0 + i * 2.4)
    bg = NAVY if i == 4 else LIGHT_GRAY
    tc = WHITE if i == 4 else NAVY
    if i > 0:
        add_shape_bg(slide, x, header_y, Inches(2.2), Inches(0.8), bg)
    add_text_box(slide, x, header_y + Inches(0.05), Inches(2.2), Inches(0.7),
                 h, 13, tc, True, PP_ALIGN.CENTER)

rows = [
    ("Framework Agnostic", "\u2717", "\u2717", "\u2713", "\u2713"),
    ("Open Source", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("AI-Native", "\u2713", "\u2713", "\u2717", "\u2713"),
    ("Lightweight", "\u2713", "\u2717", "\u2717", "\u2713"),
    ("Deterministic", "\u2717", "\u2717", "\u2713", "\u2713"),
    ("Developer-First", "\u2713", "\u2717", "\u2717", "\u2713"),
]

for r, (label, *vals) in enumerate(rows):
    y = Inches(3.2 + r * 0.65)
    bg = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(1.0), y, Inches(12), Inches(0.6), bg)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(2.0), Inches(0.5),
                 label, 14, NAVY, True)
    for c, v in enumerate(vals):
        x = Inches(3.4 + c * 2.4)
        color = TEAL if v == "\u2713" else CORAL
        if c == 3:
            color = TEAL if v == "\u2713" else CORAL
        add_text_box(slide, x, y + Inches(0.05), Inches(2.2), Inches(0.5),
                     v, 18, color, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.5),
             "No competitor is open-source, boundary-based, AND framework-agnostic.",
             16, TEAL, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: Go-to-Market
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(1.0), Inches(0.6), Inches(10), Inches(0.8),
             "Go-to-Market", 44, WHITE, True)

phases = [
    ("Months 1\u20133", "Developer\nAdoption",
     "GitHub launch\nBlog posts & tutorials\nCommunity engagement",
     "500 stars\n50 active users"),
    ("Months 4\u20136", "Ecosystem\nGrowth",
     "Framework partnerships\nIntegration guides\nCase studies",
     "2,000 stars\n200 active users"),
    ("Months 7\u201312", "Enterprise\nPipeline",
     "Launch Pro tier\nSecurity conferences\nEnterprise pilots",
     "5,000 stars\n10 pilots, $300K ARR"),
]

for i, (time, name, actions, targets) in enumerate(phases):
    x = Inches(1.0 + i * 3.8)
    y = Inches(1.8)

    # Phase number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1.2), y + Inches(0.1), Inches(0.8), Inches(0.6),
                 str(i + 1), 28, NAVY, True, PP_ALIGN.CENTER)

    add_text_box(slide, x, y + Inches(1.0), Inches(3.4), Inches(0.4),
                 time, 16, TEAL, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.4), Inches(3.4), Inches(0.6),
                 name, 20, WHITE, True, PP_ALIGN.CENTER)

    add_multi_text(slide, x + Inches(0.3), y + Inches(2.3), Inches(2.8), Inches(1.5),
                   actions.split("\n"), 14, RGBColor(0x99, 0x99, 0xAA))

    add_shape_bg(slide, x + Inches(0.3), y + Inches(4.0), Inches(2.8), Inches(1.0),
                RGBColor(0x25, 0x25, 0x3E))
    add_multi_text(slide, x + Inches(0.5), y + Inches(4.1), Inches(2.4), Inches(0.8),
                   targets.split("\n"), 14, TEAL, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: Financial Projections
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_accent_bar(slide, Inches(1.0), Inches(0.8))
add_text_box(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.8),
             "Financial Projections", 44, NAVY, True)

# Year 1 quarterly
add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5), Inches(0.5),
             "Year 1 \u2014 Quarterly Growth", 22, NAVY, True)

q_data = [
    ("Q1", "200", "0", "$0"),
    ("Q2", "500", "5", "$5K"),
    ("Q3", "800", "18", "$20K"),
    ("Q4", "1,200", "32", "$50K"),
]

# Column headers
q_headers = ["Quarter", "OSS Users", "Paid", "MRR"]
for i, h in enumerate(q_headers):
    x = Inches(1.0 + i * 1.4)
    add_shape_bg(slide, x, Inches(2.6), Inches(1.3), Inches(0.5), NAVY)
    add_text_box(slide, x, Inches(2.65), Inches(1.3), Inches(0.4),
                 h, 13, WHITE, True, PP_ALIGN.CENTER)

for r, row in enumerate(q_data):
    for c, val in enumerate(row):
        x = Inches(1.0 + c * 1.4)
        y = Inches(3.2 + r * 0.55)
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        add_shape_bg(slide, x, y, Inches(1.3), Inches(0.5), bg)
        color = TEAL if c == 3 else DARK_TEXT
        add_text_box(slide, x, y + Inches(0.05), Inches(1.3), Inches(0.4),
                     val, 14, color, c == 3, PP_ALIGN.CENTER)

# 3-year projection
add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5), Inches(0.5),
             "3-Year Trajectory", 22, NAVY, True)

years = [
    ("Year 1", "1,200 users", "$300K ARR"),
    ("Year 2", "5,000 users", "$1.5M ARR"),
    ("Year 3", "15,000 users", "$5M ARR"),
]

# Simple bar chart representation
for i, (year, users, arr) in enumerate(years):
    x = Inches(7.0)
    y = Inches(2.8 + i * 1.3)
    bar_widths = [Inches(1.5), Inches(3.0), Inches(5.0)]

    add_text_box(slide, x, y, Inches(1.5), Inches(0.4), year, 16, NAVY, True)
    bar = add_shape_bg(slide, x, y + Inches(0.4), bar_widths[i], Inches(0.5), TEAL)
    add_text_box(slide, x + bar_widths[i] + Inches(0.1), y + Inches(0.45), Inches(1.5), Inches(0.4),
                 arr, 16, NAVY, True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.45), Inches(1.5), Inches(0.4),
                 users, 12, WHITE, False)

# ============================================================
# SLIDE 11: The Ask
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text_box(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
             "The Ask", 44, WHITE, True, PP_ALIGN.CENTER)

# Funding
add_shape_bg(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.8),
            RGBColor(0x25, 0x25, 0x3E))

add_text_box(slide, Inches(2.0), Inches(2.2), Inches(9), Inches(0.5),
             "$500K \u2013 $1M Seed Round", 28, TEAL, True, PP_ALIGN.CENTER)

use_of_funds = [
    "Build and launch the open-source core (Phase 1\u20132)",
    "Hire founding security engineer",
    "Reach 1,000 active users and 10 enterprise pilots within 12 months",
]
add_multi_text(slide, Inches(2.5), Inches(2.9), Inches(8), Inches(1.0),
               [f"\u2022  {u}" for u in use_of_funds], 16,
               RGBColor(0xBB, 0xBB, 0xCC), Pt(4), alignment=PP_ALIGN.LEFT)

# Three audience boxes
audiences = [
    ("For Investors", "Own a stake in the\nopen standard for\nAI agent security"),
    ("For Early Customers", "Shape the product.\nDesign partnership.\nPriority access."),
    ("For Contributors", "Real impact.\nCommunity recognition.\nMaintainer paths."),
]

for i, (title, desc) in enumerate(audiences):
    x = Inches(1.5 + i * 3.5)
    y = Inches(4.2)
    add_shape_bg(slide, x, y, Inches(3.0), Inches(2.2), RGBColor(0x25, 0x25, 0x3E))
    add_text_box(slide, x, y + Inches(0.3), Inches(3.0), Inches(0.5),
                 title, 18, TEAL, True, PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.3), y + Inches(0.9), Inches(2.4), Inches(1.2),
                   desc.split("\n"), 14, RGBColor(0xBB, 0xBB, 0xCC),
                   alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_shape_bg(slide, Inches(0), Inches(0), W, Inches(0.05), TEAL)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "SafeAI", 72, WHITE, True, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(2.0), Inches(3.0), Inches(9), Inches(2.0), [
    "Networking got firewalls.",
    "APIs got gateways.",
    "Microservices got service meshes.",
    "AI agents are next.",
], 24, RGBColor(0x99, 0x99, 0xAA), Pt(12), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8),
             "The open standard for AI agent security.",
             22, TEAL, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
             "Let's build it together.", 18, RGBColor(0x88, 0x88, 0x99),
             False, PP_ALIGN.CENTER)

# Save
output_path = "/Users/frank.enendu/SafeAI/SafeAI-docs/SafeAI-pitch-deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
