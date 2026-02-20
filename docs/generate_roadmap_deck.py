"""Generate SafeAI Roadmap slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1A, 0x1A, 0x2E)
TEAL = RGBColor(0x16, 0xC7, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
SUBTLE = RGBColor(0x6C, 0x75, 0x7D)
DARK_CARD = RGBColor(0x25, 0x25, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tb


def multi(slide, l, t, w, h, lines, sz=16, color=SUBTLE, spacing=Pt(4), align=PP_ALIGN.LEFT, bold=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = spacing


# SLIDE 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, NAVY)
rect(slide, Inches(0), Inches(0), W, Inches(0.05), TEAL)
txt(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1), "SafeAI", 72, WHITE, True, PP_ALIGN.CENTER)
txt(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8), "Product Roadmap", 32, TEAL, False, PP_ALIGN.CENTER)
txt(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
    "Ship the smallest useful thing first. Expand based on real feedback.",
    18, RGBColor(0x88, 0x88, 0x99), False, PP_ALIGN.CENTER)

# SLIDE 2: Overview Timeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, WHITE)
rect(slide, Inches(1.0), Inches(0.8), Inches(0.6), Inches(0.06), TEAL)
txt(slide, Inches(1.0), Inches(0.95), Inches(10), Inches(0.7), "Roadmap Overview", 44, NAVY, True)

phases = [
    ("P1", "Foundation", "Wk 1\u20134", TEAL),
    ("P2", "Tool Control", "Wk 5\u20138", RGBColor(0x10, 0x80, 0x6B)),
    ("P3", "Secrets", "Wk 9\u201312", RGBColor(0x0E, 0x60, 0x50)),
    ("P4", "Proxy", "Wk 13\u201318", NAVY),
    ("P5", "Dashboard", "Wk 19\u201326", RGBColor(0x2C, 0x3E, 0x50)),
    ("P6", "Ecosystem", "Wk 27+", RGBColor(0x44, 0x44, 0x55)),
]

y_base = Inches(2.2)
# Timeline line
rect(slide, Inches(1.5), y_base + Inches(0.45), Inches(10.5), Inches(0.08), LIGHT_GRAY)

for i, (num, name, weeks, color) in enumerate(phases):
    x = Inches(1.2 + i * 1.85)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), y_base, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    txt(slide, x + Inches(0.35), y_base + Inches(0.15), Inches(0.9), Inches(0.6),
        num, 20, WHITE, True, PP_ALIGN.CENTER)
    txt(slide, x, y_base + Inches(1.1), Inches(1.6), Inches(0.5),
        name, 16, NAVY, True, PP_ALIGN.CENTER)
    txt(slide, x, y_base + Inches(1.5), Inches(1.6), Inches(0.4),
        weeks, 13, SUBTLE, False, PP_ALIGN.CENTER)

# Key milestones
txt(slide, Inches(1.0), Inches(4.5), Inches(10), Inches(0.5),
    "Key Milestones", 22, NAVY, True)

milestones = [
    ("Week 4", "PyPI Launch \u2014 SDK with input/output scanning + policy engine"),
    ("Week 8", "Tool Interception \u2014 LangChain adapter + tool contracts"),
    ("Week 12", "Secret Handling \u2014 Ephemeral credentials + approval workflows"),
    ("Week 18", "Proxy Mode \u2014 Sidecar/gateway deployment + 1000 rps"),
    ("Week 26", "Dashboard \u2014 Web UI for policy management + compliance reports"),
]

for i, (week, desc) in enumerate(milestones):
    y = Inches(5.1 + i * 0.42)
    row_bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    rect(slide, Inches(1.0), y, Inches(11), Inches(0.4), row_bg)
    txt(slide, Inches(1.2), y + Inches(0.03), Inches(1.2), Inches(0.35), week, 13, TEAL, True)
    txt(slide, Inches(2.5), y + Inches(0.03), Inches(9), Inches(0.35), desc, 13, NAVY)


# SLIDES 3-8: Individual phase slides
phase_details = [
    ("Phase 1: Foundation", "Weeks 1\u20134", "Core SDK that works out of the box",
     [("Input scanner", "Text classification with built-in PII detectors", "P0"),
      ("Output guard", "Scan responses, redact or block sensitive data", "P0"),
      ("Policy engine", "YAML-based policies with default-deny", "P0"),
      ("Data classifier", "Regex-based detection, configurable patterns", "P0"),
      ("Memory controller", "Schema-based memory, allow-listed fields", "P0"),
      ("Audit logger", "Structured JSON to stdout/file", "P0"),
      ("CLI tools", "safeai init, scan, validate", "P1"),
      ("Default policies", "Starter set covering common rules", "P1")],
     "Published to PyPI. Developer integrates in < 30 min. < 20ms overhead."),

    ("Phase 2: Tool Control", "Weeks 5\u20138", "Control what agents can do",
     [("Tool call interceptor", "Validate tool calls against contracts", "P0"),
      ("Tool contracts", "YAML declarations of tool data boundaries", "P0"),
      ("Response filtering", "Strip unauthorized fields from responses", "P0"),
      ("Agent identity", "Per-agent policy scoping", "P1"),
      ("LangChain adapter", "Framework middleware integration", "P1"),
      ("Custom classifiers", "User-defined regex in policy files", "P1"),
      ("CLI: safeai logs", "Query and display audit trail", "P1")],
     "Tool calls validated. Unauthorized data stripped. LangChain works in < 10 LOC."),

    ("Phase 3: Secrets & Approvals", "Weeks 9\u201312", "Credentials and human gates",
     [("Capability credentials", "Ephemeral secret injection with TTL", "P0"),
      ("Secret backends", "Env vars, HashiCorp Vault integration", "P0"),
      ("Human approval workflows", "CLI-based approval for high-risk actions", "P0"),
      ("Memory retention", "Automatic purge of expired data", "P1"),
      ("Encrypted handles", "Sensitive fields as encrypted references", "P1"),
      ("Claude + Google ADK", "Additional framework adapters", "P1")],
     "Secrets never in agent context. Capabilities expire. Two secret backends."),

    ("Phase 4: Proxy & Scale", "Weeks 13\u201318", "Deploy as infrastructure",
     [("HTTP proxy mode", "FastAPI-based sidecar deployment", "P0"),
      ("Gateway mode", "Centralized proxy for multi-agent environments", "P0"),
      ("Hot policy reload", "Changes without restart", "P0"),
      ("A2A enforcement", "Inter-agent trust boundaries", "P1"),
      ("Health checks", "/health endpoint with status", "P1"),
      ("Prometheus metrics", "Request counts, latency, decisions", "P1")],
     "< 50ms p99 latency. 1000+ rps. Works with any language via HTTP."),

    ("Phase 5: Dashboard & Enterprise", "Weeks 19\u201326", "Visibility and control",
     [("Web dashboard", "Overview, audit search, policy management", "P0"),
      ("Approval UI", "Web-based approval workflow", "P0"),
      ("Compliance reports", "Audit reports for time ranges", "P1"),
      ("Multi-tenant policies", "Per-team policy sets", "P1"),
      ("Role-based access", "Dashboard access control", "P1"),
      ("Alerting rules", "Configurable violation alerts", "P1")],
     "Security team investigates without engineering. Compliance reports satisfy SOC 2/GDPR."),

    ("Phase 6: Ecosystem", "Weeks 27+", "Become the standard",
     [("Plugin system", "Community classifiers and backends", "P1"),
      ("More adapters", "CrewAI, AutoGen, custom frameworks", "P1"),
      ("Structured data", "JSON/XML field-level classification", "P1"),
      ("File scanning", "Pre-process uploaded files", "P2"),
      ("Voice support", "Redaction for speech-to-text", "P2"),
      ("Browser extension", "Client-side protection", "P2")],
     "Broad framework support. Community contributions. Industry standard."),
]

for title, timing, goal, features, criteria in phase_details:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    rect(slide, Inches(1.0), Inches(0.8), Inches(0.6), Inches(0.06), TEAL)
    txt(slide, Inches(1.0), Inches(0.95), Inches(8), Inches(0.7), title, 40, NAVY, True)
    txt(slide, Inches(9.5), Inches(1.0), Inches(3), Inches(0.5), timing, 20, TEAL, True, PP_ALIGN.RIGHT)
    txt(slide, Inches(1.0), Inches(1.7), Inches(10), Inches(0.5), goal, 20, SUBTLE)

    # Feature table headers
    headers = ["Feature", "Description", "Priority"]
    widths = [Inches(2.5), Inches(6.5), Inches(1.2)]
    x_positions = [Inches(1.0), Inches(3.5), Inches(10.0)]

    y = Inches(2.5)
    for i, (h, w, x) in enumerate(zip(headers, widths, x_positions)):
        rect(slide, x, y, w, Inches(0.45), NAVY)
        txt(slide, x + Inches(0.1), y + Inches(0.05), w, Inches(0.35), h, 13, WHITE, True)

    for r, (fname, fdesc, fpri) in enumerate(features):
        y = Inches(3.05 + r * 0.48)
        row_bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        for w, x in zip(widths, x_positions):
            rect(slide, x, y, w, Inches(0.44), row_bg)
        txt(slide, x_positions[0] + Inches(0.1), y + Inches(0.05), widths[0], Inches(0.34), fname, 13, NAVY, True)
        txt(slide, x_positions[1] + Inches(0.1), y + Inches(0.05), widths[1], Inches(0.34), fdesc, 13, SUBTLE)
        pri_color = TEAL if fpri == "P0" else (RGBColor(0xF5, 0xA6, 0x23) if fpri == "P1" else SUBTLE)
        txt(slide, x_positions[2] + Inches(0.1), y + Inches(0.05), widths[2], Inches(0.34), fpri, 13, pri_color, True, PP_ALIGN.CENTER)

    # Success criteria
    rect(slide, Inches(1.0), Inches(6.5), Inches(10.2), Inches(0.6), NAVY)
    txt(slide, Inches(1.2), Inches(6.55), Inches(1.5), Inches(0.5), "Exit Criteria:", 13, TEAL, True)
    txt(slide, Inches(2.7), Inches(6.55), Inches(8.3), Inches(0.5), criteria, 13, WHITE)


# Save
path = "/Users/frank.enendu/SafeAI/SafeAI-docs/SafeAI-roadmap-deck.pptx"
prs.save(path)
print(f"Roadmap deck saved to: {path}")
print(f"Total slides: {len(prs.slides)}")
